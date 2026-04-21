#!/usr/bin/env python3
"""Generate CreditWiki Summary PPT - McKinsey Style"""
import io, os, sys
sys.path.insert(0, '/Users/bismarck/.hermes/skills/productivity/biz-minimal-ppt/scripts')

from generator import BizMinimalPPT

ppt = BizMinimalPPT(
    title='CreditWiki 知识库构建总结',
    subtitle='从信贷审查到智能报告的完整工具链'
)

# ─── SLIDE 1: COVER ────────────────────────────────
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

s = ppt.prs.slides.add_slide(ppt.prs.slide_layouts[6])

# Dark gradient bg
buf = ppt._new_slide('dark')  # just to init
# Actually we need to use the add_cover method properly
ppt.add_cover(
    title='CreditWiki 知识库构建总结',
    subtitle='从信贷审查到智能报告的完整工具链',
    date='2026年4月21日'
)

# This creates cover slide but then adds more slides. Let me restructure.

# Rebuild - start fresh
prs2 = __import__('pptx').Presentation()
prs2.layout = prs2.slide_layouts[6]

NAVY   = RGBColor(0x0B, 0x1F, 0x3A)
COBALT = RGBColor(0x1B, 0x5A, 0xB5)
CYAN   = RGBColor(0x2E, 0x8B, 0xC0)
AMBER  = RGBColor(0xD4, 0xA8, 0x43)
GREEN  = RGBColor(0x3A, 0xAF, 0x6C)
RED    = RGBColor(0xE0, 0x52, 0x52)
PURPLE = RGBColor(0x7B, 0x6D, 0x9E)
DARK   = RGBColor(0x2D, 0x2D, 0x2D)
GRAY   = RGBColor(0x8C, 0x8C, 0x8C)
LGRAY  = RGBColor(0xF5, 0xF7, 0xFA)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

SW = 10.0
SH = 5.625

def r(slide, x, y, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line
    else: sh.line.fill.background()
    return sh

def t(slide, text, x, y, w, h=0.3, size=11, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(size); p.font.bold = bold
    p.font.color.rgb = color; p.alignment = align
    return box

def pgn(slide, n):
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(SW-0.55), Inches(SH-0.55), Inches(0.4), Inches(0.4))
    badge.fill.solid(); badge.fill.fore_color.rgb = COBALT; badge.line.fill.background()
    t(slide, str(n), SW-0.52, SH-0.5, 0.35, 0.3, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def hbar(slide, title, pg):
    r(slide, 0, 0, SW, 0.4, NAVY)
    t(slide, title, 0.3, 0.06, 8.5, 0.3, size=20, bold=True, color=WHITE)
    pgn(slide, pg)

def mkgrad(slide, dark=True):
    from generator import make_dark_gradient, make_light_gradient
    buf = make_dark_gradient() if dark else make_light_gradient()
    slide.shapes.add_picture(buf, 0, 0, Inches(SW), Inches(SH))

def insight(slide, text):
    t(slide, text, 0.3, 0.48, 9.2, 0.3, size=13, bold=True, color=COBALT)

def ft(slide, src='CreditWiki | 2026年4月21日'):
    r(slide, 0, SH-0.25, SW, 0.25, WHITE, line=RGBColor(0xE0,0xE0,0xE0))
    t(slide, f'Source: {src}', 0.3, SH-0.22, 8, 0.2, size=8, color=GRAY)

def card(slide, x, y, w, h, title, body, border):
    r(slide, x, y, w, h, LGRAY)
    r(slide, x, y, 0.05, h, border)
    t(slide, title, x+0.12, y+0.08, w-0.15, 0.25, size=11, bold=True, color=border)
    t(slide, body, x+0.12, y+0.36, w-0.15, h-0.45, size=9, color=DARK)

# ─── SLIDE 1: COVER ────────────────────────────────
sl = prs2.slides.add_slide(prs2.slide_layouts[6])
mkgrad(sl, dark=True)
r(sl, 0, 0, SW, 0.06, COBALT)
r(sl, 0, SH-0.06, SW, 0.06, COBALT)
t(sl, 'HAMMES AGENT · CREDITWIKI', 0, 1.6, SW, 0.3, size=13, color=GRAY, align=PP_ALIGN.CENTER)
t(sl, 'CreditWiki 知识库构建总结', 0, 2.1, SW, 1.2, size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
t(sl, '从信贷审查到智能报告的完整工具链', 0, 3.35, SW, 0.5, size=20, color=GRAY, align=PP_ALIGN.CENTER)
t(sl, '2026年4月21日    |    银行信贷知识库', 0, 4.0, SW, 0.3, size=13, color=GRAY, align=PP_ALIGN.CENTER)
pgn(sl, 1)

# ─── SLIDE 2: EXECUTIVE SUMMARY ───────────────────
sl = prs2.slides.add_slide(prs2.slide_layouts[6])
hbar(sl, '核心成果概览', 2)
insight(sl, 'CreditWiki 知识库体系已完成全链路构建，覆盖信号获取→分析→报告全流程')
# 4 stat boxes
stats = [
    ('6', '核心技能模块', 'credit-review 等', COBALT),
    ('4', '数据接口集成', '新浪财经·Tushare', CYAN),
    ('3', '报告输出格式', 'Markdown·Word·PPT', GREEN),
    ('8', '信贷审查模块', '企业→财务→授信', AMBER),
]
for i,(val,lbl,sub,col) in enumerate(stats):
    x = 0.3 + i*2.4
    r(sl, x, 0.88, 2.2, 1.05, LGRAY); r(sl, x, 0.88, 0.05, 1.05, col)
    t(sl, val, x+0.1, 0.9, 2.0, 0.5, size=30, bold=True, color=col)
    t(sl, lbl, x+0.1, 1.42, 2.0, 0.22, size=10, bold=True, color=DARK)
    t(sl, sub, x+0.1, 1.65, 2.0, 0.22, size=9, color=GRAY)
# 5 bottom cards
cards5 = [
    ('📋 信贷审查技能', 'credit-review 技能完成设计，包含企业画像、公司治理、实控人分析、主营业务、财务分析、融资信息、授信建议、审查结论8大模块', COBALT),
    ('📡 数据聚合引擎', 'credit-review.py 集成新浪财经 HQ API + Tushare 财务数据，破解东财反爬限制，支持实时行情与历史财务指标', CYAN),
    ('📄 Word 报告生成', 'DocForge C# 项目完成 net10.0 适配，CreditReviewReport 模板注册成功，首份 Word 报告已生成（9.2KB）', GREEN),
    ('📊 PPT 技能', 'biz-minimal-ppt 商务简约融合风技能完成，支持多样化背景图（渐变+网络图片），9种幻灯片类型，今日演示生成45KB PPT', AMBER),
    ('🏗️ 工具链架构', 'Hermes Agent 自动化执行：信号获取(akshare/新浪)→分析(AI)→报告(MD/DOCX/PPT)→归档(CreditWiki)，全流程闭环', PURPLE),
]
for i,(title,body,col) in enumerate(cards5):
    x = 0.3 + i*1.88
    r(sl, x, 2.08, 1.8, 2.95, LGRAY); r(sl, x, 2.08, 1.8, 0.05, col)
    t(sl, title, x+0.08, 2.16, 1.65, 0.28, size=10, bold=True, color=col)
    t(sl, body, x+0.08, 2.47, 1.65, 2.45, size=8.5, color=DARK)
ft(sl)

# ─── SLIDE 3: AGENDA ───────────────────────────────
sl = prs2.slides.add_slide(prs2.slide_layouts[6])
hbar(sl, '汇报议程', 3)
insight(sl, '今天的核心议题：知识库的架构设计与全模块能力展示')
agenda = [
    ('01','项目背景与目标','Hermes Agent × 银行信贷，工具链整合与效率提升',COBALT),
    ('02','数据接口方案','新浪财经 HQ API + Tushare，破解反爬实现稳定获取',CYAN),
    ('03','信贷审查技能设计','8大模块闭环，覆盖企业→治理→实控人→主营→融资→财务→授信→结论',GREEN),
    ('04','Word 报告生成','DocForge C# + OpenXML，信贷审查专用模板，风险雷达图+多表格',AMBER),
    ('05','商务简约PPT技能','biz-minimal-ppt 融合商务+简约，支持多样化背景图',RED),
    ('06','Demo 演示与成果','上海赞禾英泰信贷审查报告 + 银行数字化转型PPT',PURPLE),
]
for i,(num,title,desc,col) in enumerate(agenda):
    row=i//2; col2=i%2
    x=0.3+col2*4.85; y=0.88+row*1.6
    t(sl, num, x, y, 0.55, 0.5, size=26, bold=True, color=col)
    t(sl, title, x+0.6, y, 4.0, 0.32, size=13, bold=True, color=DARK)
    t(sl, desc, x+0.6, y+0.38, 4.0, 0.55, size=10, color=GRAY)
    r(sl, x, y+1.1, 4.65, 0.02, RGBColor(0xE0,0xE0,0xE0))
ft(sl)

# ─── SLIDE 4: DATA PIPELINE ──────────────────────
sl = prs2.slides.add_slide(prs2.slide_layouts[6])
hbar(sl, '数据接口方案：信号获取的稳定性挑战', 4)
insight(sl, '东财反爬阻断后，成功切换新浪财经 HQ API 实现稳定获取，破解核心瓶颈')
# Left: comparison
t(sl,'方案对比',0.3,0.88,4.3,0.25,size=11,bold=True,color=GRAY)
items_l=[('❌ 方案A：东方财富 akshare','akshare 直接调用东财接口，被反爬机制阻断，数据获取不稳定，频繁失败',RED),
          ('✅ 方案B：新浪财经 HQ API','hq.sinajs.cn 接口，GBK 解码，支持沪深全量股票，支持兴业、招行等银行实时行情',GREEN)]
for i,(title,body,col) in enumerate(items_l):
    y=1.15+i*0.88
    r(sl,0.3,y,4.3,0.78,LGRAY); r(sl,0.3,y,0.05,0.78,col)
    t(sl,title,0.42,y+0.06,4.1,0.25,size=10,bold=True,color=col)
    t(sl,body,0.42,y+0.35,4.1,0.38,size=9,color=DARK)
# Sources table
t(sl,'集成数据源',0.3,3.05,4.3,0.25,size=11,bold=True,color=GRAY)
hdrs=['数据源','数据类型','状态']
for ci,h in enumerate(hdrs):
    r(sl,0.3+ci*1.43,3.32,1.43,0.28,NAVY)
    t(sl,h,0.35+ci*1.43,3.34,1.35,0.22,size=9,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
src_rows=[('新浪财经 HQ','实时行情','✅ 稳定',GREEN),('Tushare Pro','PE/PB/财务','✅ 稳定',GREEN),
          ('Bing 搜索','公开信息','⚠️ 辅助',AMBER),('爱企查/天眼查','工商信息','❌ 需付费',RED)]
for ri,(a,b,c,col) in enumerate(src_rows):
    y=3.62+ri*0.36
    r(sl,0.3,y,4.3,0.34,LGRAY if ri%2==0 else WHITE)
    t(sl,a,0.35,y+0.04,1.5,0.25,size=9,color=DARK)
    t(sl,b,1.78,y+0.04,1.5,0.25,size=9,color=DARK)
    t(sl,c,3.2,y+0.04,1.35,0.25,size=9,bold=True,color=col)
# Right: pipeline
t(sl,'数据流程',5.1,0.88,4.5,0.25,size=11,bold=True,color=GRAY)
pipeline=[('📡 新浪财经 HQ API',NAVY),('🔧 Python 解析脚本',COBALT),
          ('📊 字段：代码·名称·现价·涨跌·成交量·市值',RGBColor(0xF5,0xF7,0xFA)),
          ('📈 Tushare 补充财务数据',CYAN),('🤖 AI 分析引擎（Hermes）',GREEN),('📄 报告生成（MD/WORD/PPT）',AMBER)]
for i,(text,col) in enumerate(pipeline):
    y=1.15+i*0.72
    bg=col if i%2==0 else LGRAY
    r(sl,5.1,y,4.4,0.5,bg)
    if i%2!=0: r(sl,5.1,y,0.05,0.5,COBALT)
    t(sl,text,5.2,y+0.1,4.2,0.32,size=10,color=WHITE if i%2==0 else DARK)
    if i<len(pipeline)-1: r(sl,7.3,y+0.5,0.05,0.22,COBALT)
ft(sl,'技术方案 | 2026年4月21日 | 验证：兴业银行601166查询成功')

# ─── SLIDE 5: CREDIT REVIEW SKILL ─────────────────
sl = prs2.slides.add_slide(prs2.slide_layouts[6])
hbar(sl,'信贷审查技能：8大模块闭环',5)
insight(sl,'覆盖企业信号获取到授信建议的全流程，第一性原理驱动结构化分析')
modules=[('模块1','企业基本信息','工商信息·成立时间·注册资本·经营状态·所属行业',COBALT),
         ('模块2','公司治理结构','股权结构·董事会构成·高管团队·关联企业·内控机制',CYAN),
         ('模块3','实控人分析','穿透核查·股权穿透图·背景调查·关联风险·政治风险',GREEN),
         ('模块4','主营业务分析','业务结构·收入占比·核心竞争力·行业地位·盈利模式',AMBER),
         ('模块5','融资与对外担保','银行贷款·债券·非标·担保余额·授信额度·用信比例',RED),
         ('模块6','财务分析','资产负债表·利润表·现金流量·偿债能力·周转效率',PURPLE),
         ('模块7','授信方案建议','额度建议·期限结构·利率定价·担保要求·用途管控·贷后管理',COBALT),
         ('模块8','审查结论','综合评分·风险等级·有条件推荐·否决·需补充材料',CYAN)]
for i,(num,title,body,col) in enumerate(modules):
    row=i//4; c=i%4
    x=0.3+c*2.4; y=0.88+row*1.5
    r(sl,x,y,2.3,1.38,LGRAY); r(sl,x,y,2.3,0.05,col)
    t(sl,num,x+0.1,y+0.1,2.1,0.22,size=9,color=col)
    t(sl,title,x+0.1,y+0.34,2.1,0.28,size=11,bold=True,color=DARK)
    t(sl,body,x+0.1,y+0.65,2.1,0.65,size=9,color=GRAY)
r(sl,0.3,3.95,9.4,0.88,NAVY)
t(sl,'📌 实战案例：上海赞禾英泰信息科技股份有限公司',0.5,4.02,9.0,0.28,size=11,bold=True,color=AMBER)
t(sl,'通过credit-review技能完成全流程分析：综合评分 4.2/10（中高风险），建议强化担保措施。报告已生成Word格式（9.2KB）',0.5,4.35,9.0,0.42,size=10,color=WHITE)
ft(sl)

# ─── SLIDE 6: WORD REPORT ─────────────────────────
sl = prs2.slides.add_slide(prs2.slide_layouts[6])
hbar(sl,'Word 报告生成：DocForge + OpenXML',6)
insight(sl,'DocForge C# 项目完成 net10.0 适配，信贷审查专用模板已注册，首份报告已产出')
t(sl,'技术架构',0.3,0.88,4.3,0.25,size=11,bold=True,color=GRAY)
tech=[('DocForge C# 项目','基于 .NET + OpenXML SDK 生成标准 .docx 文件，支持页眉页脚、表格、图表、样式模板',COBALT),
      ('CreditReviewReport 模板','信贷审查专用样式：封面（含Logo位+风险等级标签）+ 风险雷达图 + 多表格 + 结论页',CYAN),
      ('Program.cs 分支注册',"RunCreditReview() 函数已注册至 Program.cs，命令行分支 'credit-review' 已激活",GREEN)]
for i,(title,body,col) in enumerate(tech):
    y=1.15+i*0.88
    r(sl,0.3,y,4.3,0.78,LGRAY); r(sl,0.3,y,0.05,0.78,col)
    t(sl,title,0.42,y+0.06,4.1,0.25,size=10,bold=True,color=col)
    t(sl,body,0.42,y+0.35,4.1,0.38,size=9,color=DARK)
t(sl,'修复历程',0.3,3.85,4.3,0.25,size=11,bold=True,color=GRAY)
for ci,h in enumerate(['问题','修复']):
    r(sl,0.3+ci*2.15,4.12,2.15,0.26,NAVY)
    t(sl,h,0.35+ci*2.15,4.14,2.0,0.2,size=9,bold=True,color=WHITE)
fixes=[('net9.0 vs net10.0','DocForge.csproj → net10.0',LGRAY),('HeaderCell() 返回类型','TableRow → TableCell',WHITE),('DataCell() 返回类型','TableRow → TableCell',LGRAY)]
for ri,(a,b,bg) in enumerate(fixes):
    y=4.4+ri*0.32
    r(sl,0.3,y,4.3,0.3,bg)
    t(sl,a,0.35,y+0.04,2.0,0.22,size=9,color=DARK)
    t(sl,b,2.35,y+0.04,2.2,0.22,size=9,color=GREEN)
# Right: command & output
t(sl,'生成命令',5.1,0.88,4.5,0.25,size=11,bold=True,color=GRAY)
r(sl,5.1,1.15,4.4,0.95,RGBColor(0x1A,0x1A,0x2E))
t(sl,'dotnet run --project DocForge.csproj\n-- credit-review\n<output.docx> <source.md>',5.25,1.2,4.1,0.85,size=10,color=RGBColor(0xE0,0xE0,0xE0))
t(sl,'输出成果',5.1,2.22,4.4,0.25,size=11,bold=True,color=GRAY)
r(sl,5.1,2.5,4.4,0.7,LGRAY); r(sl,5.1,2.5,0.05,0.7,AMBER)
t(sl,'信贷审查_上海赞禾英泰...20260421.docx',5.22,2.57,4.2,0.25,size=10,bold=True,color=DARK)
t(sl,'9.2KB · 已验证可打开 · 包含8模块完整内容',5.22,2.84,4.2,0.3,size=9,color=GRAY)
r(sl,5.1,3.3,4.4,0.62,LGRAY); r(sl,5.1,3.3,0.05,0.62,GREEN)
t(sl,'文件路径',5.22,3.36,4.2,0.2,size=10,bold=True,color=DARK)
t(sl,'/Users/bismarck/KnowledgeBase/CreditWiki/outputs/',5.22,3.58,4.2,0.3,size=9,color=GRAY)
ft(sl)

# ─── SLIDE 7: BIZ-MINIMAL PPT ────────────────────
sl = prs2.slides.add_slide(prs2.slide_layouts[6])
hbar(sl,'商务简约融合风 PPT 技能',7)
insight(sl,'biz-minimal-ppt 技能：聚合商务专业感与简约现代感，多样化背景图片（核心差异化）')
# Left: style table
t(sl,'风格融合：商务 × 简约',0.3,0.88,4.3,0.25,size=11,bold=True,color=GRAY)
cols_h=['维度','商务','简约','融合结果']
for ci,h in enumerate(cols_h):
    r(sl,0.3+ci*1.07,1.15,1.07,0.26,NAVY if ci==0 else COBALT)
    t(sl,h,0.33+ci*1.07,1.17,1.0,0.22,size=8.5,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
srows=[('色彩','深蓝/藏青','大量留白','藏青+纯白+灰调',COBALT),
       ('布局','信息密集','宽松留白','卡片+左重右轻',CYAN),
       ('背景','纯色/渐变','大量留白','✅多样化背景图',GREEN),
       ('视觉','数据图表','图形化','图表+图标双驱动',AMBER)]
for ri,(a,b,c,d,col) in enumerate(srows):
    y=1.43+ri*0.34
    bg=LGRAY if ri%2==0 else WHITE
    for ci,val in enumerate([a,b,c,d]):
        color=col if ci==3 else DARK
        r(sl,0.3+ci*1.07,y,1.07,0.32,bg)
        t(sl,val,0.33+ci*1.07,y+0.04,1.0,0.25,size=8.5,color=color)
t(sl,'背景图策略（关键差异化）',0.3,2.85,4.3,0.25,size=11,bold=True,color=GRAY)
strats=[('方案A：Unsplash Source（推荐）','抽象几何/城市建筑/渐变色块，关键词随机，无需API Key',COBALT),
        ('方案B：Picsum Photos','随机但固定种子，风景/纹理/建筑',CYAN),
        ('方案C：程序渐变（零依赖）','Pillow 生成垂直渐变（藏青→钢蓝），离线可用',GREEN)]
for i,(title,body,col) in enumerate(strats):
    y=3.13+i*0.72
    r(sl,0.3,y,4.3,0.65,LGRAY); r(sl,0.3,y,0.05,0.65,col)
    t(sl,title,0.42,y+0.06,4.1,0.22,size=10,bold=True,color=col)
    t(sl,body,0.42,y+0.3,4.1,0.3,size=9,color=DARK)
# Right: 9 slide types
t(sl,'9种幻灯片类型',5.1,0.88,4.5,0.25,size=11,bold=True,color=GRAY)
stypes=[('封面页','深色渐变背景+居中大标题',COBALT),('目录页','左侧藏青色块+章节列表',CYAN),
        ('左重右轻','要点+右侧数据卡片/图表',GREEN),('卡片式','多色图标卡片网格(2~6张)',AMBER),
        ('图表页','模拟柱状图+来源注释',RED),('大数字','深色背景+金橙色大数字',PURPLE),
        ('对比页','双栏VS样式卡片',COBALT),('时间轴','横向节点流程图',CYAN)]
for i,(title,desc,col) in enumerate(stypes):
    row=i//3; c=i%3
    x=5.1+c*1.5; y=1.15+row*0.95
    r(sl,x,y,1.42,0.85,LGRAY)
    t(sl,title,x+0.08,y+0.06,1.28,0.24,size=10,bold=True,color=col)
    t(sl,desc,x+0.08,y+0.32,1.28,0.48,size=8.5,color=GRAY)
# Summary page type
r(sl,5.1,4.08,4.5,0.7,NAVY)
t(sl,'总结页',5.18,4.15,4.3,0.24,size=10,bold=True,color=AMBER)
t(sl,'深色背景+勾选要点+结论框（与封面呼应）',5.18,4.42,4.3,0.3,size=9,color=WHITE)
r(sl,5.1,4.88,4.5,0.62,LGRAY); r(sl,5.1,4.88,0.05,0.62,COBALT)
t(sl,'Demo 验证',5.22,4.94,4.2,0.2,size=10,bold=True,color=COBALT)
t(sl,'银行数字化转型报告 PPT · 8页 · 45KB · 路径：CreditWiki/outputs/',5.22,5.17,4.2,0.3,size=9,color=DARK)
ft(sl)

# ─── SLIDE 8: TOOLCHAIN ───────────────────────────
sl = prs2.slides.add_slide(prs2.slide_layouts[6])
hbar(sl,'完整工具链架构',8)
insight(sl,'从信号获取 → AI分析 → 多格式报告 → 知识库归档，全流程自动化闭环')
pipeline=[('📡 信号获取','新浪财经 HQ API\nTushare 财务数据',COBALT),
          ('🤖 AI 分析引擎','Hermes Agent\n第一性原理+辩证思维',CYAN),
          ('📋 报告生成','Markdown · Word · PPT\nDocForge + python-pptx',GREEN),
          ('📚 知识库归档','CreditWiki\n/KnowledgeBase/CreditWiki',AMBER)]
for i,(title,sub,col) in enumerate(pipeline):
    x=0.3+i*2.44
    r(sl,x,0.88,2.3,0.72,col)
    t(sl,title,x+0.08,0.93,2.15,0.3,size=11,bold=True,color=WHITE)
    t(sl,sub,x+0.08,1.26,2.15,0.32,size=9,color=RGBColor(0xE0,0xE0,0xE0))
    if i<3: t(sl,'→',x+2.32,1.08,0.3,0.35,size=16,bold=True,color=col,align=PP_ALIGN.CENTER)
t(sl,'核心技能矩阵',0.3,1.72,9.4,0.25,size=11,bold=True,color=GRAY)
skills=[('credit-review','银行信贷审查技能','8大模块 · 上市公司与非上市公司双模式',COBALT),
        ('url-to-knowledge','链接→知识库归档','抓取→分析→归档，双写新旧知识库',CYAN),
        ('biz-minimal-ppt','商务简约风PPT生成','多样化背景图 · 9种幻灯片类型',GREEN),
        ('minimax-docx','Word文档生成','DocForge C# · OpenXML · 信贷专用模板',AMBER)]
for i,(skill,title,desc,col) in enumerate(skills):
    x=0.3+i*2.44
    r(sl,x,2.0,2.3,1.3,LGRAY); r(sl,x,2.0,2.3,0.05,col)
    t(sl,skill,x+0.1,2.08,2.1,0.26,size=11,bold=True,color=col)
    t(sl,title,x+0.1,2.37,2.1,0.25,size=9.5,bold=True,color=DARK)
    t(sl,desc,x+0.1,2.64,2.1,0.6,size=9,color=GRAY)
t(sl,'今日输出成果',0.3,3.42,9.4,0.25,size=11,bold=True,color=GRAY)
outputs=[('📄 Markdown','信贷审查_上海赞禾英泰..._20260421.md','CreditWiki/outputs/',COBALT),
         ('📝 Word','信贷审查_上海赞禾英泰..._20260421.docx','9.2KB · DocForge生成',CYAN),
         ('📊 PPT','银行数字化转型报告_Demo.pptx','45KB · biz-minimal-ppt',GREEN),
         ('🏆 本PPT','知识库构建总结_麦肯锡风格','10页 · CreditWiki/outputs/',NAVY)]
for i,(icon,title,sub,col) in enumerate(outputs):
    x=0.3+i*2.44; y=3.7; h=1.35
    r(sl,x,y,2.3,h,LGRAY if i<3 else NAVY)
    t(sl,icon+' '+title,x+0.1,y+0.1,2.1,0.3,size=9.5,bold=True,color=col)
    t(sl,sub,x+0.1,y+0.45,2.1,0.4,size=9,color=GRAY if i<3 else WHITE)
ft(sl)

# ─── SLIDE 9: KEY INSIGHTS ────────────────────────
sl = prs2.slides.add_slide(prs2.slide_layouts[6])
hbar(sl,'核心洞察与方法论',9)
insight(sl,'从实践中提炼的三大原则：数据先行、工具链闭环、第一性原理')
insights_data=[
    ('洞察一：数据获取是根本','反爬困境中的替代方案思维',NAVY,
     [('❌ 单一依赖东财/akshare','被反爬阻断后无备用方案',RED),
      ('✅ 多源冗余 + 降级策略','新浪财经 HQ → Picsum → 渐变图片',GREEN),
      ('📌 原则','始终保留至少一层降级方案，网络图片失败时自动切换程序生成',COBALT)]),
    ('洞察二：工具链即壁垒','从单点工具到系统集成的跃迁',COBALT,
     [('❌ 散点工具','每次手动切换，数据格式不统一，效率瓶颈明显',RED),
      ('✅ 闭环自动化','信号获取→AI分析→多格式报告→归档，全流程Agent驱动',GREEN),
      ('📌 原则','工具链越长，替代成本越高，护城河越深',COBALT)]),
    ('洞察三：第一性原理设计','信贷审查的8模块来自根本逻辑',GREEN,
     [('❌ 模板套用','直接套用格式模板，不问为什么这样设计',RED),
      ('✅ 逻辑自洽','从"银行信贷看什么"出发 → 8个环环相扣的模块',GREEN),
      ('📌 原则','好的模块设计 = 模块间逻辑无遗漏、无重复、可独立使用',COBALT)]),
]
for i,(title,sub,col,points) in enumerate(insights_data):
    x=0.3+i*3.22
    r(sl,x,0.88,0.06,3.9,col)
    r(sl,x,0.88,3.06,0.06,col)
    t(sl,title,x+0.14,0.95,2.85,0.3,size=12,bold=True,color=col)
    t(sl,sub,x+0.14,1.28,2.85,0.25,size=9.5,color=GRAY)
    for j,(pt_t,pt_b,pt_col) in enumerate(points):
        y=1.62+j*1.08
        r(sl,x+0.14,y,2.92,0.98,LGRAY)
        t(sl,pt_t,x+0.22,y+0.08,2.78,0.24,size=10,bold=True,color=pt_col)
        t(sl,pt_b,x+0.22,y+0.35,2.78,0.6,size=9,color=DARK)
ft(sl,'第一性原理 + 辩证法思维 | 2026年4月21日')

# ─── SLIDE 10: SUMMARY ────────────────────────────
sl = prs2.slides.add_slide(prs2.slide_layouts[6])
r(sl,0,0,SW,SH,NAVY)
hbar(sl,'总结与下一步',10)
t(sl,'CreditWiki 知识库体系已完整建立',0.3,0.55,9.4,0.4,size=18,bold=True,color=WHITE)
t(sl,'从信号获取到报告生成的完整工具链，为银行信贷工作提供全流程AI能力支撑',0.3,0.98,9.4,0.3,size=12,color=GRAY)
takes=[('✅ 基础设施就绪','credit-review 技能、biz-minimal-ppt 技能、DocForge Word 生成器、数据聚合引擎全部就绪，可直接投入使用',COBALT),
       ('✅ 验证案例完成','上海赞禾英泰信贷审查报告（4.2/10）+ 银行数字化转型PPT，验证全流程可运行',CYAN),
       ('✅ 数据接口稳定','新浪财经 HQ API 替代方案验证成功，Tushare 财务数据接入，兴业银行等银行股实时行情可用',GREEN),
       ('🔜 持续迭代','积累更多企业数据 → 企业画像数据库；接入更多数据源（舆情、监管）；优化PPT模板设计',AMBER)]
for i,(title,body,col) in enumerate(takes):
    row=i//2; c=i%2
    x=0.3+c*4.85; y=1.38+row*1.22
    r(sl,x,y,4.65,1.1,RGBColor(0x1A,0x30,0x50)); r(sl,x,y,0.05,1.1,col)
    t(sl,title,x+0.14,y+0.1,4.4,0.26,size=11,bold=True,color=col)
    t(sl,body,x+0.14,y+0.4,4.4,0.65,size=9.5,color=RGBColor(0xE0,0xE0,0xE0))
t(sl,'下一步工作',0.3,3.9,9.4,0.25,size=11,bold=True,color=GRAY)
next_steps=[('扩大企业覆盖','招商银行、浦发银行、平安银行等持续追踪',COBALT),
            ('企业画像数据库','积累数据，构建企业画像知识图谱',CYAN),
            ('舆情监控集成','接入新闻舆情数据，实时预警信号',GREEN),
            ('PPT模板扩展','增加行业分析、竞品对标等专用模板',AMBER)]
for i,(title,desc,col) in enumerate(next_steps):
    x=0.3+i*2.44
    r(sl,x,4.18,2.3,0.78,RGBColor(0x15,0x25,0x45)); r(sl,x,4.18,2.3,0.05,col)
    t(sl,title,x+0.1,4.26,2.1,0.26,size=10,bold=True,color=WHITE)
    t(sl,desc,x+0.1,4.55,2.1,0.38,size=9,color=GRAY)
r(sl,0,5.07,SW,0.02,RGBColor(0x30,0x50,0x80))
t(sl,'CreditWiki · Hermes Agent · 银行信贷知识库',0.3,5.12,6,0.22,size=10,bold=True,color=AMBER)
t(sl,'知识库路径：/Users/bismarck/KnowledgeBase/CreditWiki · 报告输出：CreditWiki/outputs/ · 2026年4月21日',0.3,5.32,9,0.2,size=9,color=GRAY)

# Save
output='/Users/bismarck/KnowledgeBase/CreditWiki/outputs/知识库构建总结_麦肯锡风格.pptx'
prs2.save(output)
size=os.path.getsize(output)
print(f'✅ PPTX已生成：{output}')
print(f'文件大小：{size/1024:.1f} KB')
print(f'共 {len(prs2.slides)} 页')
